#!/usr/bin/env python3
"""
PPTX Summary Generator

Generates a summary PowerPoint presentation from Design Agent output
using a PPTX template for consistent branding.

Usage:
    uvx --with python-pptx python generate-pptx.py

The uvx command (from uv package manager) handles the python-pptx dependency automatically.
"""

import os
import re
import subprocess
import sys
from pathlib import Path
from typing import Optional, List, Dict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Error: python-pptx is required. Run with: uvx --with python-pptx python generate-pptx.py")
    sys.exit(1)

# ============================================================================
# Configuration
# ============================================================================

DESIGN_DIR = Path.cwd() / "generated" / "design"
OUTPUT_DIR = Path.cwd() / "generated" / "design_pptx"
SCRIPT_DIR = Path(__file__).parent
TEMPLATE_PATH = SCRIPT_DIR.parent.parent / "templates" / "design" / "pptx-template.pptx"

# Fallback colors if template doesn't define them
COLORS = {
    "dark_blue": RGBColor(0x23, 0x2F, 0x3E),
    "orange": RGBColor(0xFF, 0x99, 0x00),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "text_dark": RGBColor(0x33, 0x33, 0x33),
}

# ============================================================================
# Content Extraction Functions
# ============================================================================

def extract_title(content: str) -> str:
    """Extract title from markdown (first # heading)"""
    match = re.search(r'^#\s+(.+?)(?:\s*[-—]|$)', content, re.MULTILINE)
    return match.group(1).strip() if match else "Design Summary"


def extract_key_value(content: str, key: str) -> Optional[str]:
    """Extract inline bold key-value: **Key**: value"""
    pattern = rf'\*\*{re.escape(key)}\*\*:\s*(.+?)(?=\n|$)'
    match = re.search(pattern, content, re.IGNORECASE)
    return match.group(1).strip() if match else None


def extract_section(content: str, heading: str) -> str:
    """Extract section content under ## Heading"""
    lines = content.split('\n')
    capturing = False
    result = []

    pattern = re.compile(rf'^##\s+{re.escape(heading)}', re.IGNORECASE)

    for line in lines:
        if pattern.match(line):
            capturing = True
            continue
        if capturing:
            if re.match(r'^##\s+', line) or re.match(r'^---\s*$', line):
                break
            result.append(line)

    return '\n'.join(result).strip()


def extract_bullet_items(content: str, limit: int = 10) -> List[str]:
    """Extract bullet list items from content"""
    items = []
    for line in content.split('\n'):
        match = re.match(r'^[-*]\s+(.+)$', line)
        if match:
            items.append(match.group(1).strip())
            if len(items) >= limit:
                break
    return items


def extract_numbered_items(content: str, limit: int = 6) -> List[Dict]:
    """Extract numbered list with bold titles: 1. **Title**: description"""
    items = []
    for line in content.split('\n'):
        match = re.match(r'^\d+\.\s+\*\*(.+?)\*\*:\s*(.+)$', line)
        if match:
            items.append({"title": match.group(1).strip(), "desc": match.group(2).strip()})
            if len(items) >= limit:
                break
    return items


def extract_key_value_bullets(content: str) -> List[Dict]:
    """Extract key-value bullet items: - **Category**: value"""
    items = []
    for line in content.split('\n'):
        match = re.match(r'^-\s+\*\*(.+?)\*\*:\s*(.+)$', line)
        if match:
            items.append({"category": match.group(1).strip(), "value": match.group(2).strip()})
    return items


def extract_mermaid_diagram(content: str) -> Optional[str]:
    """Extract the largest mermaid diagram from markdown content"""
    pattern = r'```mermaid\n([\s\S]*?)```'
    diagrams = re.findall(pattern, content)

    if not diagrams:
        return None

    # Return the diagram with the most lines (usually the main architecture diagram)
    return max(diagrams, key=lambda d: len(d.split('\n'))).strip()


def read_file_if_exists(file_path: Path) -> Optional[str]:
    """Read file if it exists"""
    try:
        if file_path.exists():
            return file_path.read_text(encoding='utf-8')
    except Exception:
        pass
    return None


# ============================================================================
# Mermaid Rendering
# ============================================================================

def render_mermaid_to_png(mermaid_code: str, output_path: Path) -> bool:
    """Render mermaid diagram to PNG using mermaid-cli (npx mmdc)"""
    mmd_path = output_path.with_suffix('.mmd')

    try:
        # Write mermaid code to .mmd file
        mmd_path.write_text(mermaid_code, encoding='utf-8')

        # Run mermaid-cli via npx
        print("  - Rendering mermaid diagram to PNG...")
        subprocess.run(
            ["npx", "mmdc", "-i", str(mmd_path), "-o", str(output_path),
             "-b", "white", "-w", "1920", "-H", "1080"],
            capture_output=True,
            timeout=120,
            check=True
        )

        if output_path.exists():
            print(f"  - Diagram rendered: {output_path}")
            return True
    except subprocess.TimeoutExpired:
        print("  - Mermaid rendering timed out")
    except subprocess.CalledProcessError as e:
        print(f"  - Mermaid rendering failed: {e}")
    except FileNotFoundError:
        print("  - npx/mmdc not found (mermaid-cli not installed)")
    except Exception as e:
        print(f"  - Mermaid rendering error: {e}")

    return False


# ============================================================================
# Template Layout Detection (AWS Template Specific)
# ============================================================================

# AWS Template layout indices (from template inspection)
LAYOUT_TITLE_SLIDE = 0           # "Title Slide" - project title slide
LAYOUT_TITLE_CONTENT = 13        # "Title and Content" - title + content placeholder
LAYOUT_TITLE_BULLETED = 14       # "Title and Bulleted Content" - title + bulleted content
LAYOUT_TWO_CONTENT = 18          # "Two Content" - title + 2 content columns
LAYOUT_COMPARISON = 19           # "Comparison" - title + 2 columns with subheadings
LAYOUT_PICTURE_CAPTION = 23      # "Picture with Caption" - title + picture + caption
LAYOUT_BLANK = 12                # "Blank" - just footer elements
LAYOUT_THANK_YOU = 41            # "Thank You" - closing slide


def get_layout_by_name(prs: Presentation, name: str):
    """Find layout by exact or partial name match"""
    name_lower = name.lower()
    for layout in prs.slide_layouts:
        if name_lower == layout.name.lower():
            return layout
    for layout in prs.slide_layouts:
        if name_lower in layout.name.lower():
            return layout
    return None


def get_layout_safe(prs: Presentation, index: int, fallback_name: str):
    """Get layout by index with fallback to name search"""
    if index < len(prs.slide_layouts):
        return prs.slide_layouts[index]
    layout = get_layout_by_name(prs, fallback_name)
    return layout if layout else prs.slide_layouts[0]


def get_content_layout(prs: Presentation):
    """Get layout for content slides (Problem, Solution)"""
    return get_layout_safe(prs, LAYOUT_TITLE_BULLETED, "Title and Bulleted Content")


def get_two_column_layout(prs: Presentation):
    """Get layout for two-column slides (Requirements, Architecture)"""
    return get_layout_safe(prs, LAYOUT_TWO_CONTENT, "Two Content")


def get_picture_layout(prs: Presentation):
    """Get layout for picture slides (Diagram)"""
    return get_layout_safe(prs, LAYOUT_PICTURE_CAPTION, "Picture with Caption")


def get_blank_layout(prs: Presentation):
    """Get blank layout"""
    return get_layout_safe(prs, LAYOUT_BLANK, "Blank")


def get_title_slide_layout(prs: Presentation):
    """Get title slide layout for project name"""
    return get_layout_safe(prs, LAYOUT_TITLE_SLIDE, "Title Slide")


def get_thank_you_layout(prs: Presentation):
    """Get thank you slide layout"""
    return get_layout_safe(prs, LAYOUT_THANK_YOU, "Thank You")


def clear_existing_slides(prs: Presentation):
    """Remove all existing slides from template (keeps layouts/masters)"""
    # Delete slides in reverse order to avoid index issues
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]


# ============================================================================
# Slide Creation Functions
# ============================================================================

def add_text_box(slide, left, top, width, height, text, font_size=9, bold=False, color=None):
    """Add a text box to a slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=8):
    """Add a bulleted list to a slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.level = 0

    return txBox


def create_title_slide(prs: Presentation, title: str, subtitle: str = "Design Specification Summary"):
    """Slide 1: Title slide with project name"""
    layout = get_title_slide_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Set the main title (placeholder 0) - match template sample positioning
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title
            shape.left = Inches(0.67)
            shape.top = Inches(2.22)
            shape.width = Inches(10.92)
        elif shape.placeholder_format.idx == 1:
            # Subtitle - position at 3.59" and match title width (10.92")
            shape.text = subtitle
            shape.left = Inches(0.67)
            shape.top = Inches(3.59)
            shape.width = Inches(10.92)

    return slide


def create_thank_you_slide(prs: Presentation):
    """Final slide: Thank you"""
    layout = get_thank_you_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Thank you layout may have placeholders for contact info, etc.
    # Leave them empty or with default template content

    return slide


def set_placeholder_text(slide, placeholder_idx: int, text: str, font_size: Optional[int] = None):
    """Set text in a placeholder by index, if it exists"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            shape.text = text
            if font_size:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size)
            return True
    return False


def fill_content_placeholder(slide, placeholder_idx: int, items: List[str], font_size: int = 18):
    """Fill a content placeholder with bullet items"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            tf = shape.text_frame
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                    p.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                p.font.size = Pt(font_size)
            return True
    return False


def create_problem_slide(prs: Presentation, data: dict):
    """Slide 2: Problem Statement"""
    layout = get_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Try to use placeholders first
    set_placeholder_text(slide, 0, "Problem Statement")  # Title placeholder

    # Build content for the body placeholder
    problem_text = data.get("problem", "Problem statement not found")
    context = f"\nTarget Delivery: {data.get('target', 'TBD')}\nScope: {data.get('scope', 'TBD')}"
    full_content = f"{problem_text}\n{context}"

    # Try placeholder with 15pt font, fall back to text box
    if not set_placeholder_text(slide, 1, full_content, font_size=15):
        add_text_box(slide, 0.5, 1.5, 12, 5, full_content, font_size=15)

    return slide


def create_solution_slide(prs: Presentation, data: dict):
    """Slide 3: Solution"""
    layout = get_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Title
    set_placeholder_text(slide, 0, "Solution")

    # Build content
    solution = data.get("solution", "Solution description not found")
    business_value = data.get("business_value", ["See executive summary"])
    bv_text = "\n".join([f"  {item}" for item in business_value])
    full_content = f"{solution}\n\nBusiness Value:\n{bv_text}"

    if not set_placeholder_text(slide, 1, full_content, font_size=15):
        add_text_box(slide, 0.5, 1.5, 12, 5, full_content, font_size=15)

    return slide


def create_functional_requirements_slide(prs: Presentation, data: dict):
    """Slide 4: Functional Requirements (two columns)"""
    layout = get_two_column_layout(prs)
    slide = prs.slides.add_slide(layout)

    func_count = data.get("func_count", 0)
    func_reqs = data.get("func_reqs", [])

    # Title
    set_placeholder_text(slide, 0, f"Functional Requirements ({func_count} total)")

    # Split into two columns
    mid = (len(func_reqs) + 1) // 2
    left_reqs = func_reqs[:mid]
    right_reqs = func_reqs[mid:]

    # Two columns with 13pt font
    if not fill_content_placeholder(slide, 1, left_reqs, font_size=13):
        add_bullet_list(slide, 0.5, 1.5, 5.5, 5.5, left_reqs, font_size=13)

    if not fill_content_placeholder(slide, 2, right_reqs, font_size=13):
        add_bullet_list(slide, 6.8, 1.5, 5.5, 5.5, right_reqs, font_size=13)

    return slide


def create_nonfunctional_requirements_slide(prs: Presentation, data: dict):
    """Slide 5: Non-Functional Requirements (two columns)"""
    layout = get_two_column_layout(prs)
    slide = prs.slides.add_slide(layout)

    nonfunc_count = data.get("nonfunc_count", 0)
    nonfunc_reqs = data.get("nonfunc_reqs", [])

    # Title
    set_placeholder_text(slide, 0, f"Non-Functional Requirements ({nonfunc_count} total)")

    # Split into two columns
    mid = (len(nonfunc_reqs) + 1) // 2
    left_reqs = nonfunc_reqs[:mid]
    right_reqs = nonfunc_reqs[mid:]

    # Two columns with 13pt font
    if not fill_content_placeholder(slide, 1, left_reqs, font_size=13):
        add_bullet_list(slide, 0.5, 1.5, 5.5, 5.5, left_reqs, font_size=13)

    if not fill_content_placeholder(slide, 2, right_reqs, font_size=13):
        add_bullet_list(slide, 6.8, 1.5, 5.5, 5.5, right_reqs, font_size=13)

    return slide


def create_architecture_slide(prs: Presentation, data: dict):
    """Slide 5: Architecture (single column)"""
    layout = get_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Title
    set_placeholder_text(slide, 0, "Architecture")

    # Prepare content - combine decisions and tech stack
    decisions = data.get("decisions", [])[:4]  # Limit to fit
    tech_stack = data.get("tech_stack", [])[:6]

    combined = []
    combined.append("Key Decisions:")
    combined.extend([f"  {d['title']}: {d['desc']}" for d in decisions])
    combined.append("")
    combined.append("Technology Stack:")
    combined.extend([f"  {t['category']}: {t['value']}" for t in tech_stack])

    # Single column content with 15pt font
    if not fill_content_placeholder(slide, 1, combined, font_size=15):
        add_bullet_list(slide, 0.67, 1.88, 12.0, 4.92, combined, font_size=15)

    return slide


def create_diagram_slide(prs: Presentation, image_path: Path):
    """Slide 6: Architecture Diagram (single column layout like slide 2)"""
    # Use single-column content layout instead of picture layout
    layout = get_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Title
    set_placeholder_text(slide, 0, "System Architecture")

    # Add image to the content area (centered, below title)
    if image_path.exists():
        # Position image - only set height (max 5"), let width scale to maintain aspect ratio
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.5), Inches(1.5),
            height=Inches(5.0)
        )

    return slide


# ============================================================================
# Main
# ============================================================================

def main():
    print("Reading design specification files...")

    # Read source files - prefer pptx-summary.md if it exists
    pptx_summary = read_file_if_exists(DESIGN_DIR / "pptx-summary.md")
    readme = read_file_if_exists(DESIGN_DIR / "README.md")
    sys_arch_file = read_file_if_exists(DESIGN_DIR / "architecture" / "system-architecture.md")

    # Architecture diagram path
    diagram_path = DESIGN_DIR / "architecture" / "system-diagram.png"

    if not pptx_summary and not readme:
        print("Error: Neither pptx-summary.md nor README.md found")
        sys.exit(1)

    # Use pptx-summary.md as primary source if available
    if pptx_summary:
        print("Using pptx-summary.md as primary source")
        source = pptx_summary

        # Extract title from README if available, otherwise from summary
        title = extract_title(readme) if readme else extract_title(source)

        # Extract from pptx-summary.md sections
        problem_section = extract_section(source, "Problem Statement")
        problem = problem_section.strip() if problem_section else "Problem statement not found"

        target_section = extract_section(source, "Target Delivery")
        target = target_section.strip() if target_section else "TBD"

        scope_section = extract_section(source, "Scope")
        scope = scope_section.strip() if scope_section else "TBD"

        solution_section = extract_section(source, "Solution")
        solution = solution_section.strip() if solution_section else "Solution description not found"

        business_value_section = extract_section(source, "Business Value")
        business_value = extract_bullet_items(business_value_section, 20)

        func_reqs_section = extract_section(source, "Functional Requirements")
        func_reqs_raw = extract_key_value_bullets(func_reqs_section)
        func_reqs = [f"{r['category']}: {r['value']}" for r in func_reqs_raw] if func_reqs_raw else extract_bullet_items(func_reqs_section, 50)

        nonfunc_reqs_section = extract_section(source, "Non-Functional Requirements")
        nonfunc_reqs_raw = extract_key_value_bullets(nonfunc_reqs_section)
        nonfunc_reqs = [f"{r['category']}: {r['value']}" for r in nonfunc_reqs_raw] if nonfunc_reqs_raw else extract_bullet_items(nonfunc_reqs_section, 50)

        decisions_section = extract_section(source, "Key Architectural Decisions")
        decisions = extract_numbered_items(decisions_section, 6)

        tech_stack_section = extract_section(source, "Technology Stack")
        tech_stack = extract_key_value_bullets(tech_stack_section)
    else:
        # Fallback to README.md and other files
        print("Falling back to README.md (pptx-summary.md not found)")
        exec_summary = read_file_if_exists(DESIGN_DIR / "project-management" / "executive-summary.md")
        func_reqs_file = read_file_if_exists(DESIGN_DIR / "requirements" / "functional-requirements.md")
        nonfunc_reqs_file = read_file_if_exists(DESIGN_DIR / "requirements" / "non-functional-requirements.md")

        title = extract_title(readme)
        problem = extract_key_value(readme, "Problem Solved") or "Problem statement not found"
        target = extract_key_value(exec_summary or readme, "Target Delivery") or "TBD"
        scope = extract_key_value(readme, "Scope") or "TBD"
        solution = extract_key_value(readme, "Solution") or "Solution description not found"

        business_value_section = extract_section(exec_summary or readme, "Business Value")
        business_value = extract_bullet_items(business_value_section, 6)

        func_reqs = extract_bullet_items(func_reqs_file or "", 10)
        nonfunc_reqs = extract_bullet_items(nonfunc_reqs_file or "", 10)

        decisions_section = extract_section(readme, "Key Architectural Decisions")
        decisions = extract_numbered_items(decisions_section, 6)

        tech_stack_section = extract_section(readme, "Technology Stack")
        tech_stack = extract_key_value_bullets(tech_stack_section)

    # Count requirements
    func_count = len(func_reqs)
    nonfunc_count = len(nonfunc_reqs)

    if not business_value:
        business_value = ["See executive summary for business value details"]

    print("Extracted content:")
    print(f"  - Problem: {len(problem)} chars")
    print(f"  - Solution: {len(solution)} chars")
    print(f"  - Business Value: {len(business_value)} items")
    print(f"  - Functional Reqs: {len(func_reqs)} items ({func_count} total)")
    print(f"  - Non-Functional Reqs: {len(nonfunc_reqs)} items ({nonfunc_count} total)")
    print(f"  - Architecture Decisions: {len(decisions)} items")
    print(f"  - Tech Stack: {len(tech_stack)} items")

    # Create output directory
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Check for architecture diagram - render if needed
    has_diagram = diagram_path.exists()

    if not has_diagram and sys_arch_file:
        mermaid_code = extract_mermaid_diagram(sys_arch_file)
        if mermaid_code:
            (DESIGN_DIR / "architecture").mkdir(parents=True, exist_ok=True)
            has_diagram = render_mermaid_to_png(mermaid_code, diagram_path)

    if has_diagram:
        print(f"  - Architecture diagram: {diagram_path}")
    elif sys_arch_file:
        print("  - No mermaid diagram found in system-architecture.md (skipping diagram slide)")
    else:
        print("  - No architecture folder (skipping diagram slide)")

    # Load template or create new presentation
    if TEMPLATE_PATH.exists():
        print(f"Using template: {TEMPLATE_PATH}")
        prs = Presentation(str(TEMPLATE_PATH))
        print(f"  - Template has {len(prs.slide_layouts)} layouts, {len(prs.slides)} sample slides")

        # Clear sample slides from template (keep masters/layouts)
        if len(prs.slides) > 0:
            print(f"  - Clearing {len(prs.slides)} sample slides from template...")
            clear_existing_slides(prs)
    else:
        print("Template not found, using default styling")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    print("\nBuilding slides...")

    # Prepare data
    data = {
        "problem": problem,
        "target": target,
        "scope": scope,
        "solution": solution,
        "business_value": business_value,
        "func_reqs": func_reqs,
        "nonfunc_reqs": nonfunc_reqs,
        "func_count": func_count,
        "nonfunc_count": nonfunc_count,
        "decisions": decisions,
        "tech_stack": tech_stack,
    }

    # Create slides in order
    create_title_slide(prs, title)                    # Slide 1: Title
    create_problem_slide(prs, data)                   # Slide 2: Problem Statement
    create_solution_slide(prs, data)                  # Slide 3: Solution
    create_functional_requirements_slide(prs, data)   # Slide 4: Functional Requirements
    create_nonfunctional_requirements_slide(prs, data) # Slide 5: Non-Functional Requirements
    create_architecture_slide(prs, data)              # Slide 6: Architecture

    if has_diagram:
        create_diagram_slide(prs, diagram_path)       # Slide 7: Diagram (optional)

    create_thank_you_slide(prs)                       # Final slide: Thank You

    # Save presentation
    output_path = OUTPUT_DIR / "design-summary.pptx"
    prs.save(str(output_path))

    slide_count = 8 if has_diagram else 7
    print(f"\nDone! Created {slide_count}-slide presentation at: {output_path}")


if __name__ == "__main__":
    main()
